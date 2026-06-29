"""
Office 文档提取模块

支持 Word (.doc, .docx) 和 PPT (.ppt, .pptx) 文档的文本提取。
优先使用 Windows COM 转换为 PDF 再用 magic-pdf 提取，
失败时回退到 python-docx / python-pptx 纯 Python 方案。
"""

import os
import platform
import shutil
import subprocess
import tempfile

from src.extractors.base import BaseExtractor
from src.logger import get_logger

logger = get_logger()


class OfficeExtractor(BaseExtractor):
    """提取 Word/PPT 文档文本。优先 COM 转换，回退纯 Python 方案。"""

    def __init__(self, config: dict):
        """
        初始化 Office 提取器。

        参数:
            config: Office 提取配置，包含 use_com, com_timeout_seconds, fallback_to_pure_python
        """
        super().__init__(config)
        self.use_com = config.get("use_com", True)
        self.com_timeout = config.get("com_timeout_seconds", 60)
        self.fallback = config.get("fallback_to_pure_python", True)

    def extract(self, file_path: str, output_dir: str) -> str | None:
        """
        提取 Office 文档文本。

        流程:
        1. 如果是 Word/PPT 且启用 COM，先尝试 COM 转 PDF，再调用 PDF 提取
        2. COM 失败或不可用，则用 python-docx/pptx 直接提取

        参数:
            file_path: Office 文件绝对路径
            output_dir: 输出目录

        返回:
            输出 .md 文件路径，失败返回 None
        """
        ext = os.path.splitext(file_path)[1].lower()
        name_without_suff = os.path.splitext(os.path.basename(file_path))[0]
        md_output = os.path.join(output_dir, f"{name_without_suff}.md")

        # 方案一：Windows COM 转换为 PDF
        if self.use_com and platform.system() == "Windows" and ext in {".doc", ".docx", ".ppt", ".pptx"}:
            try:
                pdf_path = self._convert_to_pdf_com(file_path, output_dir)
                if pdf_path:
                    # 调用 PDF 提取器
                    from src.extractors.pdf_extractor import PDFExtractor
                    pdf_extractor = PDFExtractor(self._config)
                    result = pdf_extractor.extract(pdf_path, output_dir)
                    # 清理临时 PDF
                    if pdf_path and pdf_path != file_path:
                        try:
                            os.remove(pdf_path)
                        except OSError:
                            pass
                    if result:
                        return result
            except Exception as e:
                logger.warning(f"COM 转换失败 {file_path}: {e}")

        # 方案二：纯 Python 回退
        if ext in {".doc", ".docx", ".ppt", ".pptx"}:
            pdf_path = self._convert_to_pdf_libreoffice(file_path)
            if pdf_path:
                try:
                    from src.extractors.pdf_extractor import PDFExtractor
                    pdf_extractor = PDFExtractor(self._config)
                    result = pdf_extractor.extract(pdf_path, output_dir)
                    if result:
                        return result
                finally:
                    try:
                        os.remove(pdf_path)
                    except OSError:
                        pass

        # 方案三：纯 Python 回退
        if self.fallback:
            try:
                if ext == ".docx":
                    text = self._extract_text_docx(file_path)
                elif ext == ".pptx":
                    text = self._extract_text_pptx(file_path)
                elif ext == ".doc":
                    logger.warning(f".doc 旧格式需要 COM 转换或 LibreOffice，跳过: {file_path}")
                    return None
                elif ext == ".ppt":
                    logger.warning(f".ppt 旧格式需要 COM 转换或 LibreOffice，跳过: {file_path}")
                    return None
                else:
                    return None

                if text.strip():
                    os.makedirs(output_dir, exist_ok=True)
                    with open(md_output, "w", encoding="utf-8") as f:
                        f.write(text)
                    logger.info(f"Office 提取成功 (纯Python): {file_path} → {md_output}")
                    return md_output
            except Exception as e:
                logger.error(f"纯 Python 提取失败 {file_path}: {e}")

        return None

    def _convert_to_pdf_libreoffice(self, input_path: str) -> str | None:
        """
        使用 LibreOffice headless 将 Office 文档转 PDF。

        Linux 容器中没有 Windows COM，因此优先走 LibreOffice，随后交给
        PDFExtractor/MinerU 统一提取，保留版面、公式和扫描页处理能力。
        """
        soffice = shutil.which("soffice") or shutil.which("libreoffice")
        if not soffice:
            return None

        tmp_dir = tempfile.mkdtemp(prefix="office_pdf_")
        try:
            cmd = [
                soffice,
                "--headless",
                "--nologo",
                "--nofirststartwizard",
                "--convert-to",
                "pdf",
                "--outdir",
                tmp_dir,
                os.path.abspath(input_path),
            ]
            completed = subprocess.run(
                cmd,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                text=True,
                timeout=self.com_timeout,
                check=False,
            )
            if completed.returncode != 0:
                logger.warning(f"LibreOffice 转 PDF 失败 {input_path}: {completed.stderr.strip()}")
                return None

            pdf_name = f"{os.path.splitext(os.path.basename(input_path))[0]}.pdf"
            pdf_path = os.path.join(tmp_dir, pdf_name)
            if os.path.exists(pdf_path):
                stable_path = os.path.join(tempfile.gettempdir(), f"{os.path.splitext(pdf_name)[0]}_lo.pdf")
                shutil.copy2(pdf_path, stable_path)
                return stable_path

            logger.warning(f"LibreOffice 转 PDF 完成但未找到输出: {input_path}")
            return None
        except subprocess.TimeoutExpired:
            logger.warning(f"LibreOffice 转 PDF 超时 {input_path}")
            return None
        except Exception as e:
            logger.warning(f"LibreOffice 转 PDF 出错 {input_path}: {e}")
            return None
        finally:
            shutil.rmtree(tmp_dir, ignore_errors=True)

    def supports_extension(self, ext: str) -> bool:
        """判断是否支持 Office 文件扩展名。"""
        return ext.lower() in {".doc", ".docx", ".ppt", ".pptx"}

    def _convert_to_pdf_com(self, input_path: str, output_dir: str) -> str | None:
        """
        使用 Windows COM 将 Word/PPT 转为 PDF。

        参数:
            input_path: 输入 Office 文件路径
            output_dir: 输出目录

        返回:
            生成的 PDF 文件路径，失败返回 None
        """
        try:
            from comtypes.client import CreateObject
        except ImportError:
            logger.error("comtypes 未安装，无法使用 COM 转换")
            return None

        ext = os.path.splitext(input_path)[1].lower()
        output_path = os.path.join(output_dir, f"{os.path.splitext(os.path.basename(input_path))[0]}_temp.pdf")

        try:
            if ext in {".docx", ".doc"}:
                word = CreateObject("Word.Application")
                word.Visible = False
                doc = word.Documents.Open(os.path.abspath(input_path))
                doc.SaveAs(os.path.abspath(output_path), FileFormat=17)  # 17 = PDF
                doc.Close()
                word.Quit()
            elif ext in {".pptx", ".ppt"}:
                powerpoint = CreateObject("PowerPoint.Application")
                presentation = powerpoint.Presentations.Open(os.path.abspath(input_path))
                presentation.SaveAs(os.path.abspath(output_path), 32)  # 32 = PDF
                presentation.Close()
                powerpoint.Quit()
            else:
                return None

            if os.path.exists(output_path):
                return output_path
            return None

        except Exception as e:
            logger.error(f"COM 转换出错 {input_path}: {e}")
            return None

    def _extract_text_docx(self, file_path: str) -> str:
        """
        使用 python-docx 提取 Word 文档文本。

        参数:
            file_path: .docx 文件路径

        返回:
            提取的纯文本
        """
        try:
            from docx import Document
        except ImportError:
            logger.error("python-docx 未安装，无法提取 .docx 文件")
            return ""

        doc = Document(file_path)
        paragraphs = []
        for para in doc.paragraphs:
            if para.text.strip():
                style = para.style.name.lower() if para.style else ""
                if "heading" in style:
                    level = style.replace("heading ", "").strip()
                    try:
                        prefix = "#" * int(level)
                    except ValueError:
                        prefix = "#"
                    paragraphs.append(f"{prefix} {para.text.strip()}")
                else:
                    paragraphs.append(para.text.strip())
        return "\n\n".join(paragraphs)

    def _extract_text_pptx(self, file_path: str) -> str:
        """
        使用 python-pptx 提取 PPTX 演示文稿文本。

        参数:
            file_path: .pptx 文件路径

        返回:
            提取的纯文本
        """
        try:
            from pptx import Presentation
        except ImportError:
            logger.error("python-pptx 未安装，无法提取 .pptx 文件")
            return ""

        prs = Presentation(file_path)
        slides_text = []
        for i, slide in enumerate(prs.slides, 1):
            slide_parts = [f"## 幻灯片 {i}"]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            slide_parts.append(para.text.strip())
            if len(slide_parts) > 1:
                slides_text.append("\n".join(slide_parts))
        return "\n\n".join(slides_text)
